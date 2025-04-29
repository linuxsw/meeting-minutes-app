from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation created in Part 1
prs = Presentation("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part1.pptx")

# --- Slide 4: Business Value - Time Efficiency & Cost Savings ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Business Value (1/4): Time Efficiency & Cost Savings"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Reduced Administrative Burden: Automates minute creation (saves 30-45 mins/meeting)."
p.level = 0

p = tf.add_paragraph()
p.text = "Faster Turnaround: Immediate availability of minutes."
p.level = 0

p = tf.add_paragraph()
p.text = "Resource Optimization: Frees up team members for high-value tasks."
p.level = 0

p = tf.add_paragraph()
p.text = "Quantifiable ROI: Potential to save significant administrative hours monthly."
p.level = 0

# --- Slide 5: Business Value - Enhanced Documentation Quality ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Business Value (2/4): Enhanced Documentation Quality"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Consistency: Standardized format and structure."
p.level = 0

p = tf.add_paragraph()
p.text = "Comprehensiveness: Captures all key points, decisions, and actions."
p.level = 0

p = tf.add_paragraph()
p.text = "Searchability: Easily indexed and searchable documents."
p.level = 0

p = tf.add_paragraph()
p.text = "Accuracy: Reduces human transcription errors."
p.level = 0

# --- Slide 6: Business Value - Improved Knowledge Management ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Business Value (3/4): Improved Knowledge Management"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Institutional Memory: Preserves organizational knowledge."
p.level = 0

p = tf.add_paragraph()
p.text = "Onboarding Acceleration: Quick access to historical context for new hires."
p.level = 0

p = tf.add_paragraph()
p.text = "Decision Traceability: Auditable trail of decisions."
p.level = 0

p = tf.add_paragraph()
p.text = "Cross-Team Alignment: Facilitates information sharing."
p.level = 0

# --- Slide 7: Business Value - Compliance & Governance ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Business Value (4/4): Compliance & Governance"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Meeting Documentation Standards: Helps meet organizational/regulatory requirements."
p.level = 0

p = tf.add_paragraph()
p.text = "Audit Readiness: Consistent documentation for audits."
p.level = 0

p = tf.add_paragraph()
p.text = "Process Adherence: Ensures all required elements are included."
p.level = 0

# Save the presentation (intermediate step)
prs.save("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part2.pptx")

print("Part 2 (Business Value slides) created.")

