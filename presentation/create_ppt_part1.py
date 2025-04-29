from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# --- Slide 1: Title Slide ---
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AI-Enhanced Meeting Minutes Generator"
subtitle.text = "Business Value Proposition & Feature Analysis\nProposal for Team Adoption"

# --- Slide 2: Introduction - The Challenge ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Challenge: Manual Meeting Minutes"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Creating meeting minutes manually is time-consuming and inefficient:"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Significant administrative burden (30-45 mins/meeting)"
p.level = 1

p = tf.add_paragraph()
p.text = "Delays in information distribution"
p.level = 1

p = tf.add_paragraph()
p.text = "Inconsistent quality and formatting"
p.level = 1

p = tf.add_paragraph()
p.text = "Risk of errors and omissions"
p.level = 1

p = tf.add_paragraph()
p.text = "Poor knowledge capture and retrieval"
p.level = 1

# --- Slide 3: Solution Introduction ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "The Solution: AI-Enhanced Meeting Minutes Generator"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "An automated workflow leveraging AI to process Microsoft Teams recordings and transcripts."
p.level = 0

p = tf.add_paragraph()
p.text = "Transforms raw meeting data into structured, accurate, and actionable minutes."
p.level = 0

p = tf.add_paragraph()
p.text = "Key Goals:"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Save Time & Reduce Costs"
p.level = 1

p = tf.add_paragraph()
p.text = "Improve Documentation Quality & Consistency"
p.level = 1

p = tf.add_paragraph()
p.text = "Enhance Knowledge Sharing & Accessibility"
p.level = 1

# Save the presentation (intermediate step)
prs.save("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part1.pptx")

print("Part 1 (Title and Introduction slides) created.")

