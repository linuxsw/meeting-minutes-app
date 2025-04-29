from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation created in Part 2
prs = Presentation("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part2.pptx")

# --- Slide 8: Key Features (1/4) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Features & Business Impact (1/4)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Multi-Format Input Processing"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Feature: Supports MP4+transcript, MP4-only, transcript-only."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Accommodates diverse recording practices, ensures comprehensive coverage."
p.level = 1

p = tf.add_paragraph()
p.text = "Template-Based Generation"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Feature: Specialized templates (regular, project, training, agile)."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Ensures relevant information capture for different meeting types."
p.level = 1

# --- Slide 9: Key Features (2/4) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Features & Business Impact (2/4)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "AI-Powered Content Processing"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Feature: Summarization, action item extraction, meeting-specific analysis."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Higher quality documentation, captures critical info automatically."
p.level = 1

p = tf.add_paragraph()
p.text = "Flexible AI Provider Integration"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Feature: Supports OpenAI, Together AI, Ollama (local), No AI."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Cost flexibility, business continuity, data sovereignty options."
p.level = 1

# --- Slide 10: Key Features (3/4) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Features & Business Impact (3/4)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Multiple Output Formats"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Feature: PDF, Word, HTML, plain text."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Compatibility with existing workflows and systems."
p.level = 1

p = tf.add_paragraph()
p.text = "Speaker Identification & Management"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Feature: Automatic recognition and manual correction/merging."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Accurate attribution for accountability and follow-up."
p.level = 1

# --- Slide 11: Key Features (4/4) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Features & Business Impact (4/4)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Docker Deployment Option"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Feature: Containerized deployment for easy setup."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Data sovereignty, simplified installation."
p.level = 1

p = tf.add_paragraph()
p.text = "Authentication Integration Readiness"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Feature: Designed for future AD/SSO integration."
p.level = 1

p = tf.add_paragraph()
p.text = "Impact: Secure integration into enterprise environments."
p.level = 1

# Save the presentation (intermediate step)
prs.save("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part3.pptx")

print("Part 3 (Key Features slides) created.")

