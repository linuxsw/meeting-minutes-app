from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation created in Part 4
prs = Presentation("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part4.pptx")

# --- Slide 16: Conclusion ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion: Significant Business Value"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "The AI-Enhanced Meeting Minutes Generator offers substantial benefits:"
p.level = 0

p = tf.add_paragraph()
p.text = "Significant Time Savings & Cost Reduction"
p.level = 1

p = tf.add_paragraph()
p.text = "Improved Documentation Quality & Consistency"
p.level = 1

p = tf.add_paragraph()
p.text = "Enhanced Knowledge Management & Accessibility"
p.level = 1

p = tf.add_paragraph()
p.text = "Flexible Architecture (AI Providers, Deployment)"
p.level = 1

p = tf.add_paragraph()
p.text = "Strong Potential for ROI"
p.level = 1

# --- Slide 17: Recommendation & Next Steps ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Recommendation & Next Steps"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Recommendation:"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Implement the AI-Enhanced Meeting Minutes Generator to optimize administrative processes and improve knowledge sharing."
p.level = 1

p = tf.add_paragraph()
p.text = "Next Steps:"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Initiate Preparation & Planning phase."
p.level = 1

p = tf.add_paragraph()
p.text = "Select pilot team and define scope."
p.level = 1

p = tf.add_paragraph()
p.text = "Proceed with Deployment & Configuration."
p.level = 1

# --- Slide 18: Q&A ---
slide_layout = prs.slide_layouts[5]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Add a centered text box for Q&A
left = top = width = height = Inches(1.0)
left = Inches(2.0)
width = Inches(6.0)
top = Inches(3.0)
height = Inches(1.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Q & A"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the final presentation
prs.save("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Final.pptx")

print("Part 5 (Conclusion slides) created. Final presentation saved.")

