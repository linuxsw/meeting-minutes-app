from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE # MSO_ANCHOR, MSO_THEME_COLOR are not used directly in this version

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_content_slide(prs, title_text, content_list):
    # Layout 1 is typically "Title and Content"
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title_text
    elif len(slide.placeholders) > 0 and slide.placeholders[0].has_text_frame:
         slide.placeholders[0].text_frame.text = title_text
    else: # Add a textbox if no title placeholder
        print(f"Warning: Could not set title for slide 	'{title_text}	' using standard placeholders. Adding a new textbox for title.")
        left = Inches(0.5); top = Inches(0.2); width = Inches(9); height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = title_shape.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(24)

    # Set body content
    # Placeholder index 1 is typically the main content body for "Title and Content" layout
    body_placeholder_found = False
    if len(slide.placeholders) > 1:
        try:
            # Attempt to use placeholder by common index for body
            body_shape = slide.placeholders[1]
            if body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()  # Clear existing content
                tf.word_wrap = True
                # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Optional

                for item in content_list:
                    p = tf.add_paragraph()
                    p.text = item.lstrip(" 	-*") # Remove leading hyphens/asterisks and spaces for cleaner text
                    p.font.size = Pt(14)
                    if item.strip().startswith("  -") or item.strip().startswith("    -"):
                        p.level = 1
                    elif item.strip().startswith("      -"):
                        p.level = 2
                    else:
                        p.level = 0
                body_placeholder_found = True
        except IndexError:
            pass # Placeholder index 1 doesn't exist

    if not body_placeholder_found:
        # Fallback: if placeholder[1] is not found or not suitable, add a new textbox for content
        print(f"Warning: Could not find standard body placeholder for slide 	'{title_text}	'. Adding content in a new textbox.")
        left = Inches(0.5); top = Inches(1.5); width = Inches(9.0); height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for item in content_list:
            p = tf.add_paragraph()
            p.text = item.lstrip(" 	-*") # Remove leading hyphens/asterisks and spaces
            p.font.size = Pt(14)
            if item.strip().startswith("  -") or item.strip().startswith("    -"):
                p.level = 1
            elif item.strip().startswith("      -"):
                p.level = 2
            else:
                p.level = 0

# Path to the existing presentation and the new one
existing_ppt_path = "/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Final.pptx"
updated_ppt_path = "/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Final_v3.pptx" # Versioning for clarity

# Load the existing presentation or create a new one if it doesn't exist
try:
    prs = Presentation(existing_ppt_path)
    print(f"Successfully loaded existing presentation: {existing_ppt_path}")
except Exception as e:
    print(f"Error loading presentation {existing_ppt_path}: {e}. Creating a new presentation.")
    prs = Presentation()
    # If creating new, ensure it has a title slide at least
    title_slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] # Standard subtitle placeholder for title slide
    title.text = "Meeting Minutes Generator - Enhanced"
    subtitle.text = "Project Overview and Updates"

# --- Slide: IP & Data Privacy Overview ---
slide_title_ip_privacy = "Intellectual Property & Data Privacy"
content_ip_privacy = [
    "Application prioritizes your data privacy and IP.",
    "Core Audio Processing (Transcription, Diarization):",
    "  - Uses local SpeechBrain; NO audio sent to third parties.",
    "  - Original audio/video content remains within your control.",
    "Optional AI Enhancements (Summarization, etc.):",
    "  - Text-based transcript may be sent to selected AI provider.",
    "  - Cloud Providers (OpenAI, Together AI): Users own input/output; data NOT used for training by default.",
    "  - Local AI (Ollama): Data remains within your local environment.",
    "User Control: Clear choice of AI provider, including 'No AI' option.",
    "Recommendation: For highly sensitive content, use local processing or local AI (Ollama).",
    "Refer to IP_AND_PRIVACY.md for full details."
]
add_content_slide(prs, slide_title_ip_privacy, content_ip_privacy)

# --- Slide: Security Measures & Best Practices ---
slide_title_security = "Data Security Measures & Best Practices"
content_security = [
    "Data in Transit: HTTPS used for all client-server and server-external API communication.",
    "Data at Rest (Application Server):",
    "  - Uploaded audio files are temporary and deleted after processing.",
    "  - No long-term server-side storage of full transcripts/documents by default.",
    "User Responsibility: Regularly review AI provider terms and privacy policies.",
    "Transparency: UI provides warnings when data is sent to cloud AI services.",
    "Refer to PRIVACY.md for general application privacy policy."
]
add_content_slide(prs, slide_title_security, content_security)

# --- Slide: SSO/AD Integration Strategy ---
slide_title_sso_ad = "Authentication: SSO/AD Integration Strategy"
content_sso_ad = [
    "Current Status: Guest mode for testing and initial use.",
    "Future Enhancement: Designed for enterprise authentication.",
    "Potential Integration Options:",
    "  - OAuth 2.0 / OpenID Connect (OIDC): Standard for modern SSO.",
    "  - SAML 2.0: Widely used for enterprise federation.",
    "  - LDAP / Active Directory Integration: For direct AD lookup.",
    "Benefits of Integration:",
    "  - Centralized user management and access control.",
    "  - Enhanced security and compliance with corporate policies.",
    "  - Single sign-on experience for users.",
    "Implementation Considerations:",
    "  - Requires backend changes and potentially new libraries.",
    "  - Configuration for specific Identity Providers (IdPs).",
    "  - Secure handling of tokens and sessions.",
    "Refer to docs/SSO_AD_INTEGRATION_GUIDE.md for details."
]
add_content_slide(prs, slide_title_sso_ad, content_sso_ad)

# --- Slide: Testing Strategy Overview ---
slide_title_testing_overview = "Comprehensive Testing Strategy"
content_testing_overview = [
    "Goal: Ensure reliability, functionality, and robustness of the application.",
    "Approach: Combination of automated tests covering different aspects of the system.",
    "Focus: Practical tests for core features, installation, and integration points.",
    "Documentation: Detailed strategy in TESTING_STRATEGY.md."
]
add_content_slide(prs, slide_title_testing_overview, content_testing_overview)

# --- Slide: Chosen Testing Frameworks ---
slide_title_frameworks = "Testing Frameworks & Rationale"
content_frameworks = [
    "Playwright (End-to-End & Frontend Testing):",
    "  - Why: Modern, cross-browser, auto-waits, excellent for UI interaction testing.",
    "  - Scope: User flows, UI components, API interactions from frontend.",
    "Pytest (Python Backend/Script Testing):",
    "  - Why: Mature, feature-rich, simple for unit/integration tests of Python code.",
    "  - Scope: `audio_processor.py`, `generate_pdf.py`.",
    "Why not Cucumber (for now)?",
    "  - Playwright & Pytest offer direct, efficient testing for current needs.",
    "  - BDD with Cucumber can be integrated later if project scales."
]
add_content_slide(prs, slide_title_frameworks, content_frameworks)

# --- Slide: Types of Tests Implemented ---
slide_title_test_types = "Implemented Test Categories"
content_test_types = [
    "Installation Script Tests: Validates `setup_app.sh` (dependencies, venv, PDF generation).",
    "Frontend E2E Tests (Playwright):",
    "  - Navigation, file upload UI, AI configuration UI, transcript review basics.",
    "Backend API Tests (Playwright Request Context):",
    "  - Tests API endpoints for expected responses and error handling.",
    "Python Script Tests (Pytest):",
    "  - Unit/Integration tests for `audio_processor.py` and `generate_pdf.py`.",
    "Integration Points (Conceptual/Graceful Failure):",
    "  - AI Providers: Checks UI, request formatting, mock responses.",
    "  - SSO/AD: UI elements, placeholder for future full integration tests."
]
add_content_slide(prs, slide_title_test_types, content_test_types)

# --- Slide: Running Tests ---
slide_title_running_tests = "Executing the Test Suite"
content_running_tests = [
    "Integrated into `setup_app.sh` script.",
    "  - Option `--run-tests` to execute all tests after setup.",
    "  - Installs test-specific dependencies (Playwright, Pytest).",
    "Manual Execution:",
    "  - Pytest: `source app_env/bin/activate && pytest tests/python`",
    "  - Playwright: `npx playwright test` (ensure dev server is running for E2E tests).",
    "Future: Potential integration with CI/CD pipelines for automated validation."
]
add_content_slide(prs, slide_title_running_tests, content_running_tests)

# Save the updated presentation
prs.save(updated_ppt_path)
print(f"Presentation updated with testing strategy and saved to {updated_ppt_path}")

