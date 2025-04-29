from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load the presentation created in Part 3
prs = Presentation("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part3.pptx")

# --- Slide 12: Implementation Steps Overview ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation Steps Overview"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "A phased approach to ensure successful adoption:"
p.level = 0

p = tf.add_paragraph()
p.text = "1. Preparation & Planning"
p.level = 1

p = tf.add_paragraph()
p.text = "2. Deployment & Configuration"
p.level = 1

p = tf.add_paragraph()
p.text = "3. Pilot Testing & Feedback"
p.level = 1

p = tf.add_paragraph()
p.text = "4. Training & Documentation"
p.level = 1

p = tf.add_paragraph()
p.text = "5. Rollout & Adoption"
p.level = 1

p = tf.add_paragraph()
p.text = "6. Ongoing Monitoring & Maintenance"
p.level = 1

# --- Slide 13: Implementation Steps (1/3) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation Steps (1/3)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "1. Preparation & Planning"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Identify stakeholders, gather requirements, secure resources, define pilot scope."
p.level = 1

p = tf.add_paragraph()
p.text = "2. Deployment & Configuration"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Choose deployment method (Docker recommended), set up environment, deploy app, configure AI & integrations."
p.level = 1

# --- Slide 14: Implementation Steps (2/3) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation Steps (2/3)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "3. Pilot Testing & Feedback"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Onboard pilot team, conduct pilot phase (2-4 weeks), gather feedback, iterate & refine."
p.level = 1

p = tf.add_paragraph()
p.text = "4. Training & Documentation"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Develop training materials (guides, FAQs), conduct user training, establish support channel."
p.level = 1

# --- Slide 15: Implementation Steps (3/3) ---
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Implementation Steps (3/3)"

tf = content.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "5. Rollout & Adoption"
p.level = 0
p.font.bold = True

p = tf.add_paragraph()
p.text = "Develop rollout plan (phased or full), communicate launch, monitor adoption."
p.level = 1

p = tf.add_paragraph()
p.text = "6. Ongoing Monitoring & Maintenance"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Monitor performance & costs, manage updates, provide ongoing support."
p.level = 1

# Save the presentation (intermediate step)
prs.save("/home/ubuntu/meeting-minutes-app/presentation/Meeting_Minutes_Proposal_Part4.pptx")

print("Part 4 (Implementation Steps slides) created.")

